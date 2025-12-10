import json
import os
from pathlib import Path
from typing import List, Dict

import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook


# Корень репозитория: .../my-chatgpt-telegram-bot
BASE_DIR = Path(__file__).resolve().parents[2]
# Папка с оригинальными документами
KNOWLEDGE_DIR = BASE_DIR / "knowledge" / "4lapy_docs"
# Куда сохраняем подготовленные чанки
OUTPUT_PATH = Path(__file__).resolve().parent / "docs.json"


def iter_files(root: Path) -> List[Path]:
    """Собираем все поддерживаемые файлы из knowledge/4lapy_docs."""
    exts = {".pdf", ".docx", ".pptx", ".xlsx"}
    result: List[Path] = []
    for dirpath, _, filenames in os.walk(root):
        for name in filenames:
            path = Path(dirpath) / name
            if path.suffix.lower() in exts:
                result.append(path)
    return result


def chunk_text(text: str, max_chars: int = 800, overlap: int = 100) -> List[str]:
    """Рубим длинный текст на куски для RAG."""
    text = (text or "").strip()
    if not text:
        return []

    if len(text) <= max_chars:
        return [text]

    chunks: List[str] = []
    start = 0
    length = len(text)

    while start < length:
        end = start + max_chars
        chunk = text[start:end]
        chunks.append(chunk.strip())
        if end >= length:
            break
        start = end - overlap  # небольшое перекрытие

    return chunks


def extract_from_pdf(path: Path) -> List[Dict]:
    rel = path.relative_to(BASE_DIR).as_posix()
    out: List[Dict] = []

    doc = fitz.open(path)
    try:
        for page_index in range(len(doc)):
            page = doc.load_page(page_index)
            text = page.get_text("text")
            for chunk in chunk_text(text):
                out.append(
                    {
                        "text": chunk,
                        "source": rel,
                        "source_file": path.name,
                        "page": page_index + 1,
                        "section": "",
                    }
                )
    finally:
        doc.close()

    return out


def extract_from_docx(path: Path) -> List[Dict]:
    rel = path.relative_to(BASE_DIR).as_posix()
    out: List[Dict] = []

    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    text = "\n".join(paragraphs)

    for chunk in chunk_text(text):
        out.append(
            {
                "text": chunk,
                "source": rel,
                "source_file": path.name,
                "page": None,
                "section": "",
            }
        )
    return out


def extract_from_pptx(path: Path) -> List[Dict]:
    rel = path.relative_to(BASE_DIR).as_posix()
    out: List[Dict] = []

    prs = Presentation(path)
    for slide_index, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = shape.text.strip()
                if t:
                    texts.append(t)
        if not texts:
            continue
        text = "\n".join(texts)
        for chunk in chunk_text(text):
            out.append(
                {
                    "text": chunk,
                    "source": rel,
                    "source_file": path.name,
                    "page": slide_index,
                    "section": "",
                }
            )
    return out


def extract_from_xlsx(path: Path) -> List[Dict]:
    rel = path.relative_to(BASE_DIR).as_posix()
    out: List[Dict] = []

    wb = load_workbook(path, data_only=True)
    try:
        for sheet in wb.worksheets:
            rows_text = []
            for row in sheet.iter_rows():
                values = [str(c.value).strip() for c in row if c.value not in (None, "")]
                if values:
                    rows_text.append(" | ".join(values))
            if not rows_text:
                continue
            text = f"Лист: {sheet.title}\n" + "\n".join(rows_text)
            for chunk in chunk_text(text):
                out.append(
                    {
                        "text": chunk,
                        "source": rel,
                        "source_file": path.name,
                        "page": None,
                        "section": sheet.title,
                    }
                )
    finally:
        wb.close()

    return out


def main():
    all_chunks: List[Dict] = []

    if not KNOWLEDGE_DIR.exists():
        raise SystemExit(f"Knowledge directory not found: {KNOWLEDGE_DIR}")

    files = iter_files(KNOWLEDGE_DIR)
    print(f"Found {len(files)} source files under {KNOWLEDGE_DIR}")

    for path in files:
        print(f"Processing {path}")
        suffix = path.suffix.lower()
        try:
            if suffix == ".pdf":
                chunks = extract_from_pdf(path)
            elif suffix == ".docx":
                chunks = extract_from_docx(path)
            elif suffix == ".pptx":
                chunks = extract_from_pptx(path)
            elif suffix == ".xlsx":
                chunks = extract_from_xlsx(path)
            else:
                chunks = []
        except Exception as e:
            print(f"Error while processing {path}: {e}")
            continue

        all_chunks.extend(chunks)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    with open(OUTPUT_PATH, "w", encoding="utf-8") as f:
        json.dump(all_chunks, f, ensure_ascii=False, indent=2)

    print(f"Saved {len(all_chunks)} chunks to {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
